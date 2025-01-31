import sys
from pptx import Presentation
from pyttsx3 import init
import time

def read_ppt(file_path, speed=100):  # speedパラメータを追加
    prs = Presentation(file_path)
    engine = init()
    engine.setProperty('rate', speed)  # 話速を設定

    for slide in prs.slides:
        shapes = sorted(slide.shapes, key=lambda shape: (shape.top, shape.left))
        for shape in shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    if text.strip():
                        engine.say(text)
                        engine.runAndWait()
                        time.sleep(0.5)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("コマンドラインからの呼び出し時には引数としてファイル名を指定してください")
        print("例: python read.py data/sample.pptx")
        sys.exit(1)

    ppt_file = sys.argv[1]
    read_ppt(ppt_file, speed=125)